"""Parser for PowerPoint presentations (.pptx)."""

from typing import List
from pptx import Presentation
from app.parsers.base_parser import BaseParser, ParsedChunk


class PptxParser(BaseParser):
    """Parse PowerPoint files extracting text from slides."""

    @staticmethod
    def supported_extensions() -> List[str]:
        return [".pptx", ".ppt"]

    async def parse(self, file_path: str, filename: str) -> List[ParsedChunk]:
        prs = Presentation(file_path)
        chunks = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            texts = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)

                # Extract table text
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = [cell.text.strip() for cell in row.cells]
                        texts.append(" | ".join(row_text))

            if texts:
                slide_text = f"[Slide {slide_num}]\n" + "\n".join(texts)
                chunks.append(ParsedChunk(
                    text=slide_text,
                    metadata={
                        "source": filename,
                        "file_type": "pptx",
                        "slide_number": slide_num,
                        "total_slides": len(prs.slides),
                    },
                    page_number=slide_num,
                ))

        return chunks
